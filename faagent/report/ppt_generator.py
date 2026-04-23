from pptx import Presentation

def generate_ppt(log_info, fa_results):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "FA Summary"

    content = slide.shapes.placeholders[1].text_frame
    content.text = f"SN: {log_info['sn']}\nStation: {log_info['station']}"

    for fa in fa_results:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Error: {fa['error_code']}"
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = f"""
Module: {fa['module']}
Root Cause: {fa['root_cause']}
Action: {fa['action']}
"""

    path = "FA_Report.pptx"
    prs.save(path)
    return path