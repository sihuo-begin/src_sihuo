# ──────────────────────────────────────────────
#  AR&R JSON Log Parser
#  Parses MT7 AR&R (Attribute) Jason logs
# ──────────────────────────────────────────────
import json
import logging
from pathlib import Path
import os
from datetime import datetime
from typing import Dict, Optional

import pandas as pd

logger = logging.getLogger(__name__)


def _deep_get(obj: dict, *keys, default=None):
    """Safely navigate nested dict."""
    for k in keys:
        if not isinstance(obj, dict):
            return default
        obj = obj.get(k, default)
        if obj is None:
            return default
    return obj


def _extract_sn(data: dict) -> str:
    """Extract device serial number from %device_id fields (platform+product+site+device)."""
    parts = [
        (_deep_get(data, "%device_id", "platform") or "").strip().upper(),
        (_deep_get(data, "%device_id", "product") or "").strip().upper(),
        (_deep_get(data, "%device_id", "site") or "").strip().upper(),
        (_deep_get(data, "%device_id", "device") or "").strip().upper(),
    ]
    sn = "".join(parts)
    return sn if sn else "UNKNOWN"


def _pass_fail(entry: dict) -> int:
    """Return 1 if pass, 0 if fail, from the 'status' field of a test entry."""
    status = str(_deep_get(entry, "status") or "").strip().lower()
    if status == "pass":
        return 1
    if status == "fail":
        return 0
    return -1


def parse_single_json(json_path: str) -> Optional[Dict]:
    """
    Parse one MT7 AR&R JSON log file.

    Returns a flat dict:
      - device_sn : device serial number
      - starttime : first test entry time
      - status    : PASS / FAIL (from OVERALL_TEST_RESULT entry)
      - <test_ref>: 1 (pass) or 0 (fail) for each test entry
    """
    data = None
    for enc in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            with open(json_path, "r", encoding=enc) as f:
                data = json.load(f)
            break
        except (UnicodeDecodeError, json.JSONDecodeError):
            continue
    if data is None:
        logger.warning(f"Failed to parse {json_path}: no valid encoding found")
        return None

    # Find the test array (under the first non-% test suite key)
    test_array = None
    for key, val in data.items():
        if key.startswith("%"):
            continue
        if isinstance(val, dict) and "test" in val:
            test_array = val["test"]
            break

    if not test_array:
        logger.warning(f"No test array found in {json_path}")
        return None

    record: Dict = {}

    # Store pass/fail (1/0) for each test entry by reference name
    for entry in test_array:
        if not isinstance(entry, dict):
            continue
        ref = str(_deep_get(entry, "reference") or "").strip()
        if ref:
            record[ref] = _pass_fail(entry)

    # Metadata
    record["sn"]        = _extract_sn(data)
    record["starttime"] = str(_deep_get(test_array[0], "time") or "") if test_array else ""

    # Overall status: look for OVERALL_TEST_RESULT entry
    overall_pf = -1
    for entry in test_array:
        ref = str(_deep_get(entry, "reference") or "").strip()
        if ref == "OVERALL_TEST_RESULT":
            overall_pf = _pass_fail(entry)
            break
    record["status"] = "PASS" if overall_pf == 1 else ("FAIL" if overall_pf == 0 else "UNKNOWN")

    return record


def parse_json_folder(
    folder_path: str,
    appraiser_names: Optional[Dict[str, str]] = None,
    # appraiser_names maps raw LOG_USER_NAME string -> display name (e.g. {"123": "A"})
    # NOTE: if not provided or empty, appraisers are auto-assigned A/B/C by time order
    recursive: bool = True,
) -> pd.DataFrame:
    """
    Parse all .json files in an AR&R folder into a single DataFrame.

    Each JSON file = one DUT test run.
    Columns: device_sn, sample_num, trial, appraiser, status, starttime, <test_items>

    Assignment logic (no LOG_USER_NAME dependency):
      1. Sort all JSON files by mtime (earliest first)
      2. Take up to 90 files; split into 3 groups of 30:
           Group 0 (files 0-29)  -> appraiser A
           Group 1 (files 30-59)  -> appraiser B
           Group 2 (files 60-89)  -> appraiser C
      3. Within each appraiser group: sort by mtime
         -> first 10 files = trial 1 (samples 1-10),
            next  10 files = trial 2 (samples 1-10),
            last  10 files = trial 3 (samples 1-10)
      4. Device SN sorted alphabetically -> sample 1-10
    """
    folder = Path(folder_path)
    pattern = "**/*.json" if recursive else "*.json"
    json_files = sorted(folder.glob(pattern), key=lambda p: p.stat().st_mtime)

    if not json_files:
        raise FileNotFoundError(f"No JSON files found in {folder_path}")

    logger.info(f"Found {len(json_files)} JSON files in {folder_path}")

    # ── Pre-parse all files ───────────────────────────────────────────────────
    _tmp = []
    for jf in json_files:
        rec = parse_single_json(str(jf))
        if rec:
            rec["_path"] = jf
            _tmp.append(rec)

    if not _tmp:
        raise ValueError(f"No valid records extracted from {folder_path}")

    # Cap at 90 records (3 appraisers × 3 trials × 10 samples)
    _tmp = _tmp[:90]

    # Auto-assign sample numbers from sorted unique device SNs
    all_sns = sorted(set(r.get("sn", "UNKNOWN") for r in _tmp), key=lambda x: str(x))
    sn_order = {sn: i + 1 for i, sn in enumerate(all_sns[:10])}
    logger.info(f"Auto SN->Sample mapping: {sn_order}")

    # ── Build rows with appraiser / trial / sample assignment ─────────────────
    # Split into 3 appraiser groups of up to 30 records each
    N_APPRAISERS = 3
    N_SAMPLES    = 10
    N_TRIALS     = 3
    RECS_PER_APPRAISER = N_SAMPLES * N_TRIALS  # 30

    records = []
    for group_idx in range(N_APPRAISERS):
        start = group_idx * RECS_PER_APPRAISER
        end   = start + RECS_PER_APPRAISER
        group = _tmp[start:end]
        if not group:
            continue

        # appraiser label: A, B, C (or from appraiser_names if provided)
        raw_id = f"group_{group_idx}"
        if appraiser_names and len(appraiser_names) > group_idx:
            appraiser = list(appraiser_names.values())[group_idx]
        else:
            appraiser = chr(65 + group_idx)  # A, B, C

        # Sort group by mtime to determine trial order
        group_sorted = sorted(group, key=lambda r: r["_path"].stat().st_mtime)

        for rec_idx, rec in enumerate(group_sorted):
            trial_num  = (rec_idx // N_SAMPLES) + 1   # 1, 2, 3
            sample_num = (rec_idx %  N_SAMPLES) + 1   # 1..10 repeating

            # Derive attribute: PASS=1, FAIL=0 (used by AR&R attribute agreement)
            _stat = str(rec.get("status", "UNKNOWN")).lower()
            _attr = 1 if "pass" in _stat else 0
            row = {
                "device_sn":    rec.get("sn", "UNKNOWN"),
                "sample_num":   sample_num,
                "trial":        trial_num,
                "appraiser":    appraiser,
                "inspector_id": rec.get("LOG_USER_NAME", ""),
                "starttime":    rec.get("starttime", ""),
                "status":       rec.get("status", "UNKNOWN"),
                "attribute":    _attr,
            }
            # Add pass/fail for each test entry (skip meta keys)
            for k, v in rec.items():
                if k not in {"sn", "status", "starttime", "_path"}:
                    row[k] = v
            records.append(row)

    if not records:
        raise ValueError(f"No records after assignment from {folder_path}")

    df = pd.DataFrame(records)
    df = df.sort_values(["appraiser", "trial", "sample_num"]).reset_index(drop=True)

    logger.info(f"Parsed {len(df)} rows from {len(json_files)} files")
    logger.info(f"  Unique appraisers : {sorted(df['appraiser'].unique())}")
    logger.info(f"  Unique samples    : {sorted(df['sample_num'].unique())}")
    logger.info(f"  Unique trials     : {sorted(df['trial'].unique())}")
    logger.info(f"  Status counts     :\n{df['status'].value_counts().to_string()}")

    return df


def parse_json_folder_raw(folder_path: str) -> pd.DataFrame:
    """Parse all JSON files and return the FULL GRR intermediate format
    (one row per JSON file / per measurement session, up to 90 rows).

    This matches the intermediate Excel produced by the GRR/CPK Jason parser.
    Use this as the template for the 'read from Excel' AR&R path.
    """
    folder = Path(folder_path)
    json_files = sorted(folder.glob("*.json"), key=lambda p: p.stat().st_mtime)
    if not json_files:
        raise FileNotFoundError(f"No JSON files found in {folder_path}")

    records = []
    for jf in json_files[:90]:
        rec = parse_single_json(str(jf))
        if not rec:
            continue
        row = {
            "device_sn":     rec.get("sn", "UNKNOWN"),
            "starttime":     rec.get("starttime", ""),
            "status":        rec.get("status", "UNKNOWN"),
            "appraiser_raw": rec.get("LOG_USER_NAME", ""),
            "platform":      rec.get("platform", ""),
            "product":       rec.get("product", ""),
            "site":          rec.get("site", ""),
            "device":        rec.get("device", ""),
        }
        for k, v in rec.items():
            if k not in {"sn", "status", "starttime", "platform",
                          "product", "site", "device", "LOG_USER_NAME", "_path"}:
                row[k] = v
        records.append(row)

    if not records:
        raise ValueError(f"No valid records from {folder_path}")

    df = pd.DataFrame(records)
    df = df.sort_values(["device_sn", "starttime"]).reset_index(drop=True)

    # Assign trial: split each SN's records into groups of 3
    def _assign_trial(g):
        g = g.copy()
        n = len(g)
        trial = [1]*min(3,n) + [2]*max(0,min(3,n-3)) + [3]*max(0,n-6)
        g["trial"] = trial[:n]
        return g
    df = df.groupby("device_sn", group_keys=False).apply(_assign_trial, include_groups=False)

    # Add pass/fail flag from OVERALL_TEST_RESULT
    if "OVERALL_TEST_RESULT" in df.columns:
        df["_pf"] = df["OVERALL_TEST_RESULT"].apply(
            lambda v: "PASS" if str(v).strip() in ("0","pass","PASS","") else "FAIL"
        )
        df.insert(3, "status_pf", df.pop("_pf"))

    logger.info(f"parse_json_folder_raw: {len(df)} rows x {len(df.columns)} cols")
    return df


# ─────────────────────────────────────────────────────────────────────────────
# Minitab AR&R (Attribute Gauge R&R) via ATGAUGE command
# ─────────────────────────────────────────────────────────────────────────────
import subprocess as _sub
import time as _time

def _run_arr_minitab(mtb_path: str, df_jason: pd.DataFrame,
                     minitab_opts: dict,
                     excel_path: str = None,
                     sample_map: dict = None) -> dict:
    """
    Generate AR&R charts via Minitab ATGAUGE command.

    Workflow:
      1. Run ATGAUGE + XPPOINT via subprocess
      2. Wait for Minitab to finish (PowerPoint opens with charts)
      3. Use GetActiveObject to get the open PowerPoint and SaveAs to target path
      4. Wait 15s to ensure all images are loaded into the presentation
      5. Extract images from PPTX using python-pptx
      6. Delete PPTX and extracted images
    """
    out_dir = Path.home() / 'arr_charts'
    out_dir.mkdir(exist_ok=True)
    ts = int(_time.time() * 1000)
    pptx_path = str(out_dir / 'arr_charts_{}.pptx'.format(ts))
    img_dir = str(out_dir)

    # Kill existing Minitab
    for _ in range(3):
        try:
            _sub.run(['taskkill', '/F', '/IM', 'Mtb.exe'],
                     capture_output=True, timeout=5)
        except Exception:
            pass
    _time.sleep(1.0)

    # Prepare data
    df = df_jason.copy().dropna(subset=['appraiser', 'sample_num', 'trial', 'status'])
    df = df.sort_values(['appraiser', 'trial', 'sample_num']).reset_index(drop=True)

    resp = df['status'].apply(
        lambda v: 1 if str(v).strip().lower() in ('pass', '1') else 0).tolist()
    sample = df['sample_num'].astype(int).tolist()
    inspector_numbers = minitab_opts.get('inspector_numbers', [])
    logger.info("_run_arr_minitab: inspector_numbers=%s", inspector_numbers)
    # Priority: (1) inspector_numbers from config mapping > (2) inspector_id from JSON > (3) raw label
    if inspector_numbers:
        uniq_apps = sorted(df['appraiser'].unique().tolist(), key=str)
        app_map = {app: inspector_numbers[i] if i < len(inspector_numbers) else str(app)
                   for i, app in enumerate(uniq_apps)}
        appra = df['appraiser'].map(app_map).astype(str).tolist()
    else:
        appra = df['appraiser'].astype(str).tolist()
    # Build reference (standard) values: from sample_map expected values
    if sample_map:
        ref = [sample_map.get(int(s), {}).get("expected", 1) for s in df["sample_num"]]
    else:
        ref = [1] * len(df)

    date_val = minitab_opts.get('date')
    date_str = date_val.strip() if (date_val and date_val.strip()) else datetime.now().strftime('%b %d %Y')
    user_str = minitab_opts.get('user') or 'Simon Huo'
    prod_str = minitab_opts.get('product', 'Alpha')

    def fmt(vals, width=10):
        return [' '.join(str(v) for v in vals[i:i+width]) for i in range(0, len(vals), width)]

    lines = (
        [
            'NAME C1 "Column"',
            'NAME C2 "Samples"',
            'NAME C3 "Appraisers"',
            'NAME C4 "attribute"',
            '',
            'SET C1',
        ] + fmt(resp) + ['END.', '',
        'SET C2',
        ] + fmt(sample) + ['END.', '',
        'SET C3',
        ] + fmt(appra) + ['END.', '',
        'SET C4',
        ] + fmt(ref) + ['END.', '',
        '',
        'ATGAUGE;',
        '  Response \'Column\';',
        '  Samples \'Samples\';',
        '  Appraisers \'Appraisers\';',
        '  Attribute \'attribute\';',
        '  Date "{}";'.format(date_str),
        '  User "{}";'.format(user_str),
        '  Product "{}";'.format(prod_str),
        '  GAPPraiser;',
        '  GAttribute;',
        '  ONEDoc;',
        '  BRIEF 2.',
        '',
        'XPPOINT.',
        '',
        'END',
    ])

    mtb_file = out_dir / 'arr_{}.mtb'.format(ts)
    for attempt in range(4):
        try:
            mtb_file.write_text('\n'.join(lines), encoding='utf-8')
            break
        except (PermissionError, OSError):
            if attempt < 3:
                _time.sleep(0.5)
            else:
                raise

    logger.info('_run_arr_minitab MTB (%d lines): %s', len(lines), mtb_file)
    logger.info('PPTX target: %s', pptx_path)

    paths = {}
    try:
        # Run Minitab: ATGAUGE creates charts, XPPOINT opens PowerPoint
        proc = _sub.Popen([mtb_path, str(mtb_file)],
                          stdout=_sub.PIPE, stderr=_sub.PIPE)
        logger.info('Minitab started (PID=%s)', proc.pid)

        # Wait for Minitab to finish running commands
        # (it exits after XPPOINT has opened PowerPoint)
        logger.info('Waiting for Minitab to finish ATGAUGE+XPPOINT...')
        try:
            proc.wait(timeout=25)
            logger.info('Minitab finished (exit code %s)', proc.returncode)
        except _sub.TimeoutExpired:
            logger.warning('Minitab did not exit in 60s, terminating...')
            proc.terminate()
            proc.wait(timeout=10)

        # Give PowerPoint a moment to fully load all charts
        logger.info('Waiting 5s for PowerPoint to initialise...')
        _time.sleep(5)

        # Use GetActiveObject to get the open PowerPoint and SaveAs
        try:
            import win32com.client
            import os as _os

            logger.info('Connecting to open PowerPoint via GetActiveObject...')
            ppt_app = win32com.client.GetActiveObject('Powerpoint.Application')
            presentation = ppt_app.ActivePresentation

            _os.makedirs(out_dir, exist_ok=True)
            # ppSaveAsOpenXMLPresentation = 24
            presentation.SaveAs(pptx_path, FileFormat=24)
            logger.info('PPTX saved via COM: %s', pptx_path)

            # Wait 15s to ensure all images are fully loaded into the presentation
            logger.info('Waiting 15s for all images to load into PPTX...')
            _time.sleep(3)

            # Close PowerPoint cleanly
            try:
                presentation.Saved = True
                ppt_app.Quit()
                logger.info('PowerPoint closed')
            except Exception as e:
                logger.warning('PowerPoint.Quit failed: %s', e)

            # Release COM objects so PowerPoint fully exits
            del presentation
            del ppt_app
            import gc
            gc.collect()
            _time.sleep(2.0)
            logger.info('COM objects released')

        except ImportError:
            logger.warning('win32com not available, PPTX save skipped')
        except Exception as e:
            logger.warning('PowerPoint COM SaveAs failed: %s', e)

        # Kill Minitab
        try:
            _sub.run(['taskkill', '/F', '/IM', 'Mtb.exe'],
                     capture_output=True, timeout=5)
        except Exception:
            pass

        # Extract images from PPTX
        # Wait for PPTX file to be fully released by PowerPoint after Quit
        for _wait_idx in range(6):
            if Path(pptx_path).exists():
                break
            logger.info('Waiting for PPTX to be released (%d)...', _wait_idx)
            _time.sleep(1.0)
        if Path(pptx_path).exists():
            logger.info('PPTX confirmed: %s (size=%d)', pptx_path, Path(pptx_path).stat().st_size)
            try:
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                prs = Presentation(pptx_path)
                ext_map = {'image/jpeg': '.jpg', 'image/png': '.png',
                           'image/gif': '.gif', 'image/bmp': '.bmp'}

                extracted = []
                for slide_idx, slide in enumerate(prs.slides, 1):
                    for shape_idx, shape in enumerate(slide.shapes, 1):
                        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                            continue
                        img = shape.image
                        ext = ext_map.get(img.ext, '.png')
                        out_name = 'slide{:02d}_shape{:02d}{}'.format(
                            slide_idx, shape_idx, ext)
                        out_path = _os.path.join(img_dir, out_name)
                        with open(out_path, 'wb') as f:
                            f.write(img.blob)
                        extracted.append(out_path)
                        logger.info('Extracted: %s', out_path)

                if extracted:
                    paths['arr_main'] = extracted[0]
                    if len(extracted) >= 2:
                        paths['arr_gage'] = extracted[1]
                    if len(extracted) > 2:
                        for i, p in enumerate(extracted[2:], 3):
                            paths['arr_chart_{}'.format(i)] = p

                logger.info('Extracted %d images from PPTX', len(extracted))

                try:
                    Path(pptx_path).unlink()
                    logger.info('PPTX deleted')
                except Exception as e:
                    logger.warning('Could not delete PPTX: %s', e)

                paths['_img_dir'] = img_dir
                paths['_extracted'] = extracted

            except ImportError:
                logger.warning('python-pptx not installed')
            except Exception as e:
                logger.warning('PPTX extraction failed: %s', e)
        else:
            logger.warning('PPTX not found: %s', pptx_path)

    except Exception as e:
        logger.warning('_run_arr_minitab failed: %s', e)
        try:
            _sub.run(['taskkill', '/F', '/IM', 'Mtb.exe'],
                     capture_output=True, timeout=5)
        except Exception:
            pass

    return paths
