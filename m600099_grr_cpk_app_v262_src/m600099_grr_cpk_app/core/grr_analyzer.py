# ──────────────────────────────────────────────
#  GRR Analyzer  –  AIAG (Average & Range) method
#
#  Model:
#    Total TV   = √(EV² + PV²)          (Total Variation)
#    EV         = Repeatability  = d₂·R̄ᵢ  (Equipment Variation)
#    PV         = Reproducibility= d₂·XP̄   (Part Variation)
#    GRR%       = 100 × TV / (TV + PV)   [%GRR of tolerance – NIST style]
#    P/T Ratio  = 6·TV / tolerance
#    NDC        = √(PV/EV) ≈ 1.41·(PV/R) – Number of Distinct Categories
#
#  GRR < 10%  → Excellent
#  10–30%    → Acceptable
#  > 30%     → Marginal / Unacceptable
# ──────────────────────────────────────────────
import logging
import os
import re
import subprocess
import tempfile
from pathlib import Path

import numpy as np
import pandas as pd
import json
from itertools import islice, cycle
import easyocr
import time
reader = easyocr.Reader(['en'], gpu=False)

from utils.config import grr_tolerance, CPK_SPECS, LED_MAPPING, _norm_led, GRR_OPS, GRR_FORM_DEFAULTS, OUTPUT_DIR

logger = logging.getLogger(__name__)

# ── d₂ table (partial – most common values) ──
_D2 = {
    (2, 2): 1.128, (2, 3): 1.693, (2, 4): 2.059, (2, 5): 2.326,
    (2, 6): 2.534, (2, 7): 2.704, (2, 8): 2.847, (2, 9): 2.970,
    (2,10): 3.078,
    (3, 2): 1.693, (3, 3): 2.394, (3, 4): 2.772, (3, 5): 3.078,
    (3, 6): 3.258, (3, 7): 3.407, (3, 8): 3.532, (3, 9): 3.640,
    (3,10): 3.735,
}

def _d2(n, k):
    """d₂ for k repetitions, n operators / part groups."""
    return _D2.get((n, k), 1.0)   # fallback


def _parse_minitab_session_for_grr_str(text):
    """
    Parse Minitab's Gage evaluation table from a session-output text
    string (captured via `OUTFILE` or `JOURNAL`). Returns (pct, pt,
    ndc) where each is a float or None. The relevant lines are:

        |------- 95% CI --------|  |------- 95% CI --------|
    Source         VarComp  %VarComp  StDev    %SV      6*SD     SV/Toler
    Total Gage R&R <VarComp>  <%VarComp>  <StDev>  <%SV>  (<CI>)  <6*SD>  (<CI>)  <SV/Toler>  (<CI>)
    ...
    Number of Distinct Categories = <N>
    """
    import re as _re
    if not text:
        return None
    pct_val = None
    pt_val  = None
    ndc_val = None
    for raw in text.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        norm = _re.sub(r"[ \t]+", " ", line.strip())
        low  = norm.lower()

        if ndc_val is None and "number of distinct categories" in low and "=" in norm:
            try:
                tail = norm.split("=", 1)[1].strip()
                m = _re.search(r"-?\d+(?:\.\d+)?", tail)
                if m:
                    ndc_val = float(m.group(0))
            except Exception:
                pass
            continue

        if (pct_val is None or pt_val is None) \
                and "total gage r&r" in low \
                and "total gage r&r variance" not in low:
            # OCR often dumps the whole Gage Evaluation table on a
            # single line — slice the substring from "total gage
            # r&r" to the next source name.
            _m = _re.search(
                r"total gage r&r\b(.*?)(?:repeatability|"
                r"reproducibility|operator|part-to-part|"
                r"part-to_part|total variation|\Z)",
                low)
            if _m:
                _sub = _m.group(1)
                _sub = _re.sub(r"\([^)]*\)", " ", _sub)
                _nums = _re.findall(
                    r"-?\d+(?:\.\d+)?(?:[eE][-+]?\d+)?", _sub)
                if len(_nums) >= 6:
                    pct_val = float(_nums[3])
                    pt_val  = float(_nums[5])
                elif len(_nums) >= 4:
                    pct_val = float(_nums[2])
                    pt_val  = float(_nums[3])
    return (pct_val, pt_val, ndc_val)


def _preprocess_for_ocr(pil_img):
    """
    Improve OCR accuracy on Minitab chart PNGs.

    Steps:
      1. Greyscale (drop colour noise).
      2. Upscale 3× (Minitab's chart text is small at 96 dpi; 3×
         is the sweet spot — tesseract's LSTM needs ~30px cap
         height for reliable recognition).
      3. Otsu binarisation (auto-threshold to pure B/W).
      4. Median blur 1px to remove antialiasing specks.
    Returns the processed PIL.Image.
    """
    try:
        from PIL import ImageOps, ImageFilter
        img = pil_img.convert("L")
        # Upscale 3× with bicubic
        w, h = img.size
        if max(w, h) < 2400:
            img = img.resize((w * 3, h * 3), 3)  # 3 = BICUBIC
        # Otsu binarisation (auto threshold)
        try:
            import numpy as _np
            arr = _np.asarray(img)
            # Otsu: minimise intra-class variance
            hist, _ = _np.histogram(arr, bins=256, range=(0, 256))
            total = arr.size
            if total == 0:
                return img
            sum_all = _np.dot(_np.arange(256), hist)
            sum_b = 0.0
            w_b = 0.0
            max_var = -1.0
            threshold = 127
            for t in range(256):
                w_b += hist[t]
                if w_b == 0:
                    continue
                w_f = total - w_b
                if w_f == 0:
                    break
                sum_b += t * hist[t]
                m_b = sum_b / w_b
                m_f = (sum_all - sum_b) / w_f
                var = w_b * w_f * (m_b - m_f) ** 2
                if var > max_var:
                    max_var = var
                    threshold = t
            img = img.point(lambda p: 255 if p > threshold else 0)
        except Exception:
            # Fall back to simple threshold
            img = img.point(lambda p: 255 if p > 160 else 0)
        # Smooth 1px to remove isolated specks
        try:
            img = img.filter(ImageFilter.MedianFilter(3))
        except Exception:
            pass
        return img
    except Exception as pe:
        logger.debug("preprocess_for_ocr failed: %s", pe)
        return pil_img


# def _ocr_extract_grr_metrics(image_paths):
#     """
#     Run Tesseract on each chart image and parse out the
#         Total Gage R&R   %SV  SV/Toler
#     line and the
#         Number of Distinct Categories = N
#     line.
#
#     `image_paths` is a dict whose values are disk paths OR BytesIO
#     buffers (the V224 in-memory cleanup variant). Values that are
#     None or whose keys start with '_' (e.g. '_img_dir') are skipped.
#
#     Returns (pct, pt, ndc) where each is a float or None. Returns
#     (None, None, None) if no OCR backend is installed or no chart
#     image was recognised.
#     """
#     import re as _re
#     # V252: pytesseract path removed (tesseract.exe cannot be
#     # installed on this machine). Only the easyocr backend is
#     # used from here on.
#     # have_easyocr = False
#     # try:
#     #     import easyocr  # noqa: F401
#     #     have_easyocr = True
#     # except Exception:
#     #     pass
#     # if not have_easyocr:
#     #     logger.warning(
#     #         "No OCR backend available. easyocr is the only "
#     #         "supported backend now (V252+). Install with:\n"
#     #         "  pip install easyocr\n"
#     #         "First run downloads ~100MB model to "
#     #         "%%USERPROFILE%%\\.EasyOCR\\model\\")
#     #     return (None, None, None)
#
#     pct_val = None
#     pt_val  = None
#     ndc_val = None
#
#     for k, v in (image_paths or {}).items():
#         if k.startswith("_") or v is None:
#             continue
#         # Only OCR the two slides that contain the Minitab
#         # "Gage Evaluation" table and the "Number of Distinct
#         # Categories" line — everything else (header, ANOVA
#         # table, misclassification probabilities, report cover)
#         # is irrelevant for our metrics and slow to OCR. Minitab
#         # places them on slide02 shapes 3 and 4.
#         if not (("slide02_shape03" in k) or ("slide02_shape04" in k)):
#             continue
#         text = ""
#         try:
#             if hasattr(v, "read"):
#                 try:
#                     from PIL import Image as _PIL
#                     v.seek(0)
#                     img_obj = _PIL.open(v)
#                 except ImportError:
#                     logger.warning("Pillow not installed; cannot read "
#                                    "BytesIO image for OCR")
#                     continue
#             else:
#                 try:
#                     from PIL import Image as _PIL
#                     img_obj = _PIL.open(v)
#                 except ImportError:
#                     logger.warning("Pillow not installed")
#                     continue
#             # Preprocess: greyscale + upscale + Otsu binarisation
#             proc_img = _preprocess_for_ocr(img_obj)
#             # V252: easyocr only (pytesseract removed).
#             try:
#                 import tempfile as _tmp
#                 import os as _os
#                 import cv2
#                 # import easyocr as _eo
#                 _fd, _p = _tmp.mkstemp(suffix=".png")
#                 try:
#                     _os.close(_fd)
#                     proc_img.save(_p)
#                     if not hasattr(_ocr_extract_grr_metrics,
#                                    "_easyocr_reader"):
#                         _ocr_extract_grr_metrics._easyocr_reader = \
#                             reader.Reader(["en"], gpu=False, verbose=False)
#                             # reader.readtext(["en"], gpu=False, verbose=False)
#                     result = _ocr_extract_grr_metrics._easyocr_reader\
#                         .readtext(_p, detail=0, paragraph=False)
#                     text = "\n".join(result)
#                 finally:
#                     try: _os.unlink(_p)
#                     except Exception: pass
#             except Exception as ee2:
#                 logger.warning(
#                     "OCR of %s failed (easyocr): %s", k, ee2)
#                 continue
#         except Exception as oe:
#             logger.warning("OCR of %s failed: %s", k, oe)
#             continue
#         if not text:
#             continue
#         # V253: normalise the OCR text so the parser can split on
#         # the actual '|' delimiter Minitab uses. easyocr sometimes
#         # returns multi-line text where Minitab's pipe-delimited
#         # columns end up wrapped onto separate physical lines —
#         # replacing the newline with '|' makes the parser robust
#         # to that layout.
#         text = text.replace("\n", "|")
#         logger.info("OCR of image %s (first 400 chars): %s",
#                     k, text[:400].replace("\n", " | "))
#         logger.info("simon \n {}".format(text))
#         logger.info("simon line \n {}".format(text.splitlines))
#         for raw in text.splitlines():
#             norm = _re.sub(r"[ \t]+", " ", raw.strip())
#             low  = norm.lower()
#             if ndc_val is None and "number of distinct categories" in low and "=" in norm:
#                 try:
#                     tail = norm.split("=", 1)[1].strip()
#                     m = _re.search(r"-?\d+(?:\.\d+)?", tail)
#                     logger.info(
#                         "V250 OCR parser: NDC line found, tail=[%r], "
#                         "match=%s", tail, m.group(0) if m else None)
#                     if m:
#                         ndc_val = float(m.group(0))
#                 except Exception as ne:
#                     logger.warning("V250 OCR NDC parse failed: %s", ne)
#                 continue
#             if (pct_val is None or pt_val is None) \
#                     and "total gage r&r" in low \
#                     and "total gage r&r variance" not in low:
#                 # OCR often dumps the whole Gage Evaluation table
#                 # on a single line separated by '|'. V250's regex
#                 # approach (r"total gage r&r\b(.*?)(?:repeatability|...)")
#                 # was observed to return an empty sub-string in some
#                 # Minitab 22 deployments (the non-greedy .*? was
#                 # zero-matching). V251 instead splits the line on
#                 # the actual '|' delimiter Minitab uses and walks
#                 # the segments in order — no regex state machine
#                 # involved.
#                 try:
#                     _segs = [s.strip() for s in norm.split("|")]
#                     _segs_low = [s.lower() for s in _segs]
#                     logger.info("simon \n {}".format(_segs_low))
#                     # _hit = -1
#                     # for _i, _seg in enumerate(_segs_low):
#                     #     if "total gage r&r" in _seg:
#                     #         _hit = _i
#                     #         break
#                     # if _hit >= 0:
#                     #     _nums = []
#                     #     for _seg in _segs[_hit + 1:]:
#                     #         for _m in _re.findall(
#                     #                 r"-?\d+(?:\.\d+)?(?:[eE][-+]?\d+)?",
#                     #                 _seg):
#                     #             _nums.append(float(_m))
#                     #             if len(_nums) >= 4:
#                     #                 break
#                     #         if len(_nums) >= 4:
#                     #             break
#                     #     logger.info(
#                     #         "V251 OCR parser: split-by-pipe found "
#                     #         "Total Gage R&R at segment %d, "
#                     #         "next-nums=%s", _hit, _nums)
#                     #     if len(_nums) >= 4:
#                     #         # Gage Evaluation: [StDev, 6*SD, %SV, SV/Toler]
#                     #         pct_val = _nums[2]
#                     #         pt_val  = _nums[3]
#                     #     elif len(_nums) >= 6:
#                     #         pct_val = _nums[3]
#                     #         pt_val  = _nums[5]
#                     # else:
#                     #     logger.warning(
#                     #         "V251 OCR parser: 'total gage r&r' substring "
#                     #         "found in line but no '|'-segment starts with "
#                     #         "it. Line was: %r", low[:200])
#                     for key in range(len(_segs_low)):
#                         if "total gage r&r" in _segs_low[key]:
#                             pct_val = _segs_low[key+3]
#                             logger.info(',' in pct_val)
#                             if ',' in pct_val:
#                                 pct_val=pct_val.replace(',', '')
#                             pt_val  = _segs_low[key+4]
#                             logger.info(',' in pt_val)
#                             if ',' in pt_val:
#                                 pt_val=pt_val.replace(',', '')
#                             break
#                 except Exception as pe:
#                     logger.warning("V251 OCR parser exception: %s", pe)
#     if pct_val is not None or pt_val is not None or ndc_val is not None:
#         logger.info(
#             "OCR extracted: pct=%s pt=%s ndc=%s", pct_val, pt_val, ndc_val)
#     return (pct_val, pt_val, ndc_val)


def _ocr_extract_grr_metrics(image_paths):
    """
    Run Tesseract on each chart image and parse out the
        Total Gage R&R   %SV  SV/Toler
    line and the
        Number of Distinct Categories = N
    line.

    `image_paths` is a dict whose values are disk paths OR BytesIO
    buffers (the V224 in-memory cleanup variant). Values that are
    None or whose keys start with '_' (e.g. '_img_dir') are skipped.

    Returns (pct, pt, ndc) where each is a float or None. Returns
    (None, None, None) if no OCR backend is installed or no chart
    image was recognised.
    """
    import re as _re
    # V252: pytesseract path removed (tesseract.exe cannot be
    # installed on this machine). Only the easyocr backend is
    # used from here on.
    # have_easyocr = False
    # try:
    #     import easyocr  # noqa: F401
    #     have_easyocr = True
    # except Exception:
    #     pass
    # if not have_easyocr:
    #     logger.warning(
    #         "No OCR backend available. easyocr is the only "
    #         "supported backend now (V252+). Install with:\n"
    #         "  pip install easyocr\n"
    #         "First run downloads ~100MB model to "
    #         "%%USERPROFILE%%\\.EasyOCR\\model\\")
    #     return (None, None, None)

    pct_val = None
    pt_val  = None
    ndc_val = None
    for k, v in (image_paths or {}).items():
        if k.startswith("_") or v is None:
            continue
        # Only OCR the two slides that contain the Minitab
        # "Gage Evaluation" table and the "Number of Distinct
        # Categories" line — everything else (header, ANOVA
        # table, misclassification probabilities, report cover)
        # is irrelevant for our metrics and slow to OCR. Minitab
        # places them on slide02 shapes 3 and 4.
        if not (("slide02_shape03" in k) or ("slide02_shape04" in k)):
            continue
        text = ""
        try:
            if hasattr(v, "read"):

                try:
                    from PIL import Image as _PIL
                    v.seek(0)
                    img_obj = _PIL.open(v)
                except ImportError:
                    logger.warning("Pillow not installed; cannot read "
                                   "BytesIO image for OCR")
                    continue
            else:
                try:
                    from PIL import Image as _PIL
                    img_obj = _PIL.open(v)
                except ImportError:
                    logger.warning("Pillow not installed")
                    continue
            # Preprocess: greyscale + upscale + Otsu binarisation
            proc_img = _preprocess_for_ocr(img_obj)
            # V252: easyocr only (pytesseract removed).
            try:
                import tempfile as _tmp
                import os as _os
                import cv2
                # import easyocr as _eo
                _fd, _p = _tmp.mkstemp(suffix=".png")
                try:
                    _os.close(_fd)
                    proc_img.save(_p)
                    # if not hasattr(_ocr_extract_grr_metrics,
                    #                "_easyocr_reader"):
                    #     _ocr_extract_grr_metrics._easyocr_reader = \
                    #         reader.Reader(["en"], gpu=False, verbose=False)
                    #         # reader.readtext(["en"], gpu=False, verbose=False)
                    # result = _ocr_extract_grr_metrics._easyocr_reader\
                    #     .readtext(_p, detail=0, paragraph=False)
                    # text = "\n".join(result)
                    img = cv2.imread(_p)
                    h, w = img.shape[:2]
                    print(h, w)
                    # ✅ 裁剪右上角
                    # roi = img[0:int(h), int(w):w]

                    # ✅ 降采样
                    roi = cv2.resize(img, None, fx=0.6, fy=0.6)

                    # ✅ 灰度
                    gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)

                    # ✅ OCR
                    text = reader.readtext(gray)
                finally:
                    logger.info("\n final")
                    try: _os.unlink(_p)
                    except Exception: pass
            except Exception as ee2:
                logger.warning(
                    "OCR of %s failed (easyocr): %s", k, ee2)
                continue
        except Exception as oe:
            logger.warning("OCR of %s failed: %s", k, oe)
            continue
        if not text:
            continue
        # V253: normalise the OCR text so the parser can split on
        # the actual '|' delimiter Minitab uses. easyocr sometimes
        # returns multi-line text where Minitab's pipe-delimited
        # columns end up wrapped onto separate physical lines —
        # replacing the newline with '|' makes the parser robust
        # to that layout.
        # text = text.replace("\n", "|")
        logger.info("simon text is \n {}".format(str(text).split(',')))
        if " 'Total Gage R&R'" in str(text).split(','):
            for key in range(len(text)):
                if 'Total Gage R&R' in text[key]:
                    logger.info("simon text key +1 is \n {}".format(text[key + 1][1]))
                    logger.info("simon text key +1 is \n {}".format(len(text[key + 1][1].split(" "))))
                    if len(text[key + 1][1].split(" ")) > 1:
                        pct_val = text[key + 2][1]
                        pt_val = text[key + 3][1]
                    else:
                        pct_val = text[key + 3][1]
                        pt_val = text[key + 4][1]
                    logger.info(',' in pct_val)
                    if ',' or ',' in pct_val:
                        pct_val = pct_val.replace(',', '')
                        pct_val = pct_val.replace(',', '')
                    logger.info(',' in pt_val)
                    if ',' or ',' in pt_val:
                        pt_val = pt_val.replace(',', '')
                        pt_val = pt_val.replace(',', '')
                    break
        if " 'Number of Distinct Categories'" in str(text).split(','):
            for key in range(len(text)):
                if 'Number of Distinct Categories' in text[key]:
                    ndc_val = text[key + 1][1]
                    ndc_val = ndc_val.replace("=", "")

        # logger.info("OCR of image %s (first 400 chars): %s",
        #             k, text[:400].replace("\n", " | "))
        # logger.info("simon \n {}".format(text))
        # logger.info("simon line \n {}".format(text.splitlines))
        # for raw in text.splitlines():
        #     norm = _re.sub(r"[ \t]+", " ", raw.strip())
        #     low  = norm.lower()
        #     if ndc_val is None and "number of distinct categories" in low and "=" in norm:
        #         try:
        #             tail = norm.split("=", 1)[1].strip()
        #             m = _re.search(r"-?\d+(?:\.\d+)?", tail)
        #             logger.info(
        #                 "V250 OCR parser: NDC line found, tail=[%r], "
        #                 "match=%s", tail, m.group(0) if m else None)
        #             if m:
        #                 ndc_val = float(m.group(0))
        #         except Exception as ne:
        #             logger.warning("V250 OCR NDC parse failed: %s", ne)
        #         continue
        #     if (pct_val is None or pt_val is None) \
        #             and "total gage r&r" in low \
        #             and "total gage r&r variance" not in low:
        #         # OCR often dumps the whole Gage Evaluation table
        #         # on a single line separated by '|'. V250's regex
        #         # approach (r"total gage r&r\b(.*?)(?:repeatability|...)")
        #         # was observed to return an empty sub-string in some
        #         # Minitab 22 deployments (the non-greedy .*? was
        #         # zero-matching). V251 instead splits the line on
        #         # the actual '|' delimiter Minitab uses and walks
        #         # the segments in order — no regex state machine
        #         # involved.
        #         try:
        #             _segs = [s.strip() for s in norm.split("|")]
        #             _segs_low = [s.lower() for s in _segs]
        #             logger.info("simon \n {}".format(_segs_low))
        #             # _hit = -1
        #             # for _i, _seg in enumerate(_segs_low):
        #             #     if "total gage r&r" in _seg:
        #             #         _hit = _i
        #             #         break
        #             # if _hit >= 0:
        #             #     _nums = []
        #             #     for _seg in _segs[_hit + 1:]:
        #             #         for _m in _re.findall(
        #             #                 r"-?\d+(?:\.\d+)?(?:[eE][-+]?\d+)?",
        #             #                 _seg):
        #             #             _nums.append(float(_m))
        #             #             if len(_nums) >= 4:
        #             #                 break
        #             #         if len(_nums) >= 4:
        #             #             break
        #             #     logger.info(
        #             #         "V251 OCR parser: split-by-pipe found "
        #             #         "Total Gage R&R at segment %d, "
        #             #         "next-nums=%s", _hit, _nums)
        #             #     if len(_nums) >= 4:
        #             #         # Gage Evaluation: [StDev, 6*SD, %SV, SV/Toler]
        #             #         pct_val = _nums[2]
        #             #         pt_val  = _nums[3]
        #             #     elif len(_nums) >= 6:
        #             #         pct_val = _nums[3]
        #             #         pt_val  = _nums[5]
        #             # else:
        #             #     logger.warning(
        #             #         "V251 OCR parser: 'total gage r&r' substring "
        #             #         "found in line but no '|'-segment starts with "
        #             #         "it. Line was: %r", low[:200])
        #             for key in range(len(_segs_low)):
        #                 if "total gage r&r" in _segs_low[key]:
        #                     pct_val = _segs_low[key+3]
        #                     logger.info(',' in pct_val)
        #                     if ',' in pct_val:
        #                         pct_val=pct_val.replace(',', '')
        #                     pt_val  = _segs_low[key+4]
        #                     logger.info(',' in pt_val)
        #                     if ',' in pt_val:
        #                         pt_val=pt_val.replace(',', '')
        #                     break
        #         except Exception as pe:
        #             logger.warning("V251 OCR parser exception: %s", pe)
    if pct_val is not None or pt_val is not None or ndc_val is not None:
        logger.info(
            "OCR extracted: pct=%s pt=%s ndc=%s", pct_val, pt_val, ndc_val)
    return (pct_val, pt_val, ndc_val)



def augment_grr_with_ocr(grr_result, chart_paths,
                         fallback_pct=None, fallback_pt=None, fallback_ndc=None):
    """
    Run OCR on chart images and overwrite grr_result.grr_pct / pt_ratio
    / ndc with values pulled from the Minitab Gage Evaluation table.

    chart_paths may be a dict whose values are disk paths OR BytesIO
    buffers (the V224 in-memory cleanup variant). Any key whose value
    is None or whose name starts with '_' (e.g. '_img_dir') is skipped.

    On OCR failure (tesseract missing / unparseable output) the
    Python-computed fallback_* values are kept.
    """
    metrics = _ocr_extract_grr_metrics(chart_paths)
    pct_m, pt_m, ndc_m = metrics
    if pct_m is not None:
        grr_result.grr_pct = float(pct_m)
    elif fallback_pct is not None:
        grr_result.grr_pct = float(fallback_pct)
    if pt_m is not None:
        grr_result.pt_ratio = float(pt_m)
    elif fallback_pt is not None:
        grr_result.pt_ratio = float(fallback_pt)
    if ndc_m is not None:
        grr_result.ndc = float(ndc_m)
    elif fallback_ndc is not None:
        grr_result.ndc = float(fallback_ndc)
    logger.info(
        "GRR %s: OCR override pct=%s pt=%s ndc=%s (Python fallback %s/%s/%s)",
        grr_result.item, pct_m, pt_m, ndc_m,
        fallback_pct, fallback_pt, fallback_ndc)
    return grr_result


class GRRResult:
    """Holds GRR computation results for one LED item."""
    def __init__(self, item: str):
        self.item = item
        self.ev    = None   # Equipment Variation
        self.pv    = None   # Part Variation
        self.tv    = None   # Total Variation
        self.grr_pct = None # %GRR
        self.pt_ratio = None # P/T ratio
        self.ndc  = None   # Number of Distinct Categories
        self.tolerance = None
        self.summary = ""   # Text summary
        self.chart_paths = {}  # {chart_type: file_path}

    def grade(self) -> str:
        if self.grr_pct is None:
            return "N/A"
        g = self.grr_pct
        if g < 10:
            return "Excellent ✅"
        if g < 30:
            return "Acceptable ⚠️"
        return "Marginal ❌"

    def to_dict(self) -> dict:
        return {
            "item": self.item,
            "EV": round(self.ev, 4) if self.ev else None,
            "PV": round(self.pv, 4) if self.pv else None,
            "TV": round(self.tv, 4) if self.tv else None,
            "GRR%": round(self.grr_pct, 2) if self.grr_pct else None,
            "P/T": round(self.pt_ratio, 4) if self.pt_ratio else None,
            "NDC": round(self.ndc, 1) if self.ndc else None,
            "tolerance": round(self.tolerance, 2) if self.tolerance else None,
            "grade": self.grade(),
            "summary": self.summary,
            "chart_paths": self.chart_paths,
        }


class GRRAnalyzer:
    """
    AIAG Average-and-Range GRR for one LED intensity column.

    Expects df to have:
      - 'sn'      : device serial number (part identifier)
      - 'QR_SCAN' : operator / appraiser identifier (optional)
      - <item>    : measurement value
    """

    # Factor d₂ for n=2 operators, k repetitions
    D2_2_2 = 1.128

    def __init__(self, df: pd.DataFrame, item: str, LED_SPECS):
        self.df    = df
        self.item  = item
        self.result = GRRResult(item)
        self.LED_SPECS= LED_SPECS
        logger.info("simon LED_SPECS is {}".format(LED_SPECS))

    def _detect_cols(self, df: pd.DataFrame, col: str):
        """Detect sn/appraiser column names regardless of format."""
        # sn: prefer "part_num" (GRR part number, set by assign_trials) > "sn" > "Sample"
        # In intermediate Excel: sn=UNKNOWN (not useful), part_num=1..10 (correct for GRR)
        sn_candidates = [c for c in df.columns if c.lower() in ("sn", "sample", "part_num")]
        # Prefer part_num > sn > Sample (part_num is set by assign_trials for GRR)
        sn_col = None
        for candidate in ["part_num", "sn", "sample"]:
            for c in sn_candidates:
                if c.lower() == candidate:
                    sn_col = c
                    break
            if sn_col:
                break
        # Operator column: appraiser > inspector > qr_scan
        # appraiser = operator ID from MT7 QR_SCAN result (intermediate Excel)
        # inspector = GRR template format
        # qr_scan = raw MT7 QR_SCAN field (not ideal for GRR)
        op_candidates = [c for c in df.columns if c.lower() in ("appraiser", "inspector", "qr_scan", "operator")]
        op_col = op_candidates[0] if op_candidates else None
        return sn_col, op_col

    def compute(self, minitab_path: str = None) -> GRRResult:
        # Resolve actual column name in df (item may be bare LED name or PNUM)
        item = self.item
        if item in self.df.columns:
            col = item
        elif _norm_led(item) in self.df.columns:
            col = _norm_led(item)
        else:
            # Try PNUM reverse lookup
            col = None
            for pnum, led in LED_MAPPING.items():
                if led == item and pnum in self.df.columns:
                    col = pnum
                    break
            if col is None:
                logger.warning("Column for item '%s' not found in df. Available: %s", item, list(self.df.columns))
                return self.result

        sn_col, op_col = self._detect_cols(self.df, col)
        if sn_col is None:
            logger.warning("No sn/Part column found for col=%s. Available cols: %s", col, list(self.df.columns))
            return self.result

        # Build keep_cols — always include sn_col, op_col, and the LED measurement col
        if op_col:
            keep_cols = [sn_col, op_col, col]
        else:
            keep_cols = [sn_col, col]

        logger.info("compute: item=%s col=%s sn_col=%s op_col=%s keep_cols=%s df_rows=%d",
                     item, col, sn_col, op_col, keep_cols, len(self.df))

        df = self.df[keep_cols].copy()
        df["part"]       = df[sn_col].astype(str).str.strip()
        df["appraiser"]  = df[op_col].astype(str).str.strip() if op_col else "A"
        df[col]           = pd.to_numeric(df[col], errors="coerce")
        df.dropna(subset=[col], inplace=True)

        logger.info("compute: after dropna rows=%d unique_parts=%s",
                     len(df), df["part"].unique().tolist() if len(df) > 0 else "N/A")

        n_parts     = df["part"].nunique()
        n_appraisers = df["appraiser"].nunique()
        n_reps      = max(df.groupby(["part", "appraiser"]).size().max(), 2)

        logger.info(f"GRR {col}: {n_parts} parts × {n_appraisers} appraisers × {n_reps} reps")

        # ── Average & Range per appraiser ──
        X_bar = df.groupby("appraiser")[col].mean()          # X̄ per appraiser
        R_i   = df.groupby("appraiser")[col].apply(
            lambda x: x.max() - x.min()
        )                                                       # Range per appraiser
        X_bar_bar = X_bar.mean()                               # grand mean

        # ── Equipment Variation (EV) ──
        R_bar_i = R_i.mean()                                   # average range
        EV = _d2(2, n_reps) * R_bar_i / 1.128 * self.D2_2_2
        EV = _d2(2, 2) * R_i.mean()                           # simplified: d₂(2,2)=1.128

        # ── Part Variation (PV) ──
        X_p_bar = df.groupby("part")[col].mean()
        R_p     = X_p_bar.max() - X_p_bar.min()
        PV = _d2(2, 2) * R_p

        # ── Total Variation ──
        TV = np.sqrt(EV**2 + PV**2) if (EV and PV) else (EV or PV or 0)

        # ── Tolerance (USL – LSL) ──
        lo, hi = CPK_SPECS.get(col, (None, None))
        if lo is None or hi is None:
            # Fall back to 20% of mean
            tol = X_bar_bar * 0.20
        else:
            tol = hi - lo
        tolerance = tol

        # ── GRR metrics ──
        # Python-computed fallback values (used if OCR can't read the
        # official values from the Minitab chart images).
        grr_pct_py  = 100 * TV / tolerance if tolerance else None
        pt_ratio_py = 6   * TV / tolerance if tolerance else None
        ndc_py      = 1.41 * PV / EV if (PV and EV) else None

        self.result.ev          = float(EV)  if EV  else None
        self.result.pv          = float(PV)  if PV  else None
        self.result.tv          = float(TV)  if TV  else None
        # Default to Python-computed values
        self.result.grr_pct     = float(grr_pct_py)  if grr_pct_py  else None
        self.result.pt_ratio    = float(pt_ratio_py) if pt_ratio_py else None
        self.result.ndc         = float(ndc_py)      if ndc_py      else None
        self.result.tolerance   = float(tolerance)

        # ── Minitab charts (always run if a path is given) ──
        if minitab_path and os.path.isfile(minitab_path):
            chart_paths = self._run_minitab(minitab_path, df, col, timeout=30, LED_SPECS=self.LED_SPECS)
            self.result.chart_paths = chart_paths
        elif not self.result.chart_paths:
            # Batch flow (V224+): main_window injects in-memory BytesIO
            # chart paths AFTER the Minitab run. OCR will run after
            # main_window's per-item augment.
            pass

        # ── Pull %GR&R / %P/T / NDC from Minitab's actual output ──
        # Primary path: read the Minitab `PRINT M1-M30` worksheet
        # output (captured via MSAVE — the only headless way to get
        # the official Gage Evaluation numbers on Minitab 22).
        # Secondary path: OCR the chart PNGs. Last resort: Python
        # fallback.
        if minitab_path and os.path.isfile(minitab_path):
            worksheet_text = None
            try:
                # Reconstruct the worksheet file path that
                # _run_minitab wrote (same naming pattern: out_dir /
                # grr_<safe>_<ts>_vals.txt). Read the newest matching
                # file in the GRR chart dir, then delete it.
                _chart_dir = Path.home() / "grr_charts"
                if _chart_dir.exists():
                    _safe = _norm_led(col).replace("/", "_")
                    for _vp in sorted(
                            _chart_dir.glob("grr_*_vals.txt"),
                            key=lambda p: p.stat().st_mtime,
                            reverse=True):
                        if _safe in _vp.name:
                            try:
                                worksheet_text = _vp.read_text(
                                    encoding="latin1", errors="ignore")
                                logger.info(
                                    "GRR %s: read worksheet print %d chars from %s",
                                    col, len(worksheet_text), _vp.name)
                            except Exception as se:
                                logger.warning(
                                    "GRR %s: could not read worksheet print %s: %s",
                                    col, _vp.name, se)
                            try:
                                _vp.unlink(missing_ok=True)
                            except Exception:
                                pass
                            break
            except Exception as se:
                logger.warning("GRR %s: worksheet read error: %s", col, se)

            metrics = _parse_minitab_session_for_grr_str(worksheet_text)
            if metrics is not None and (metrics[0] is not None
                                        or metrics[1] is not None
                                        or metrics[2] is not None):
                pct_m, pt_m, ndc_m = metrics
                logger.info(
                    "GRR %s: JOURNAL values pct=%s pt=%s ndc=%s "
                    "(overrides Python pct=%s pt=%s ndc=%s)",
                    col, pct_m, pt_m, ndc_m,
                    grr_pct_py, pt_ratio_py, ndc_py)
                if pct_m is not None:
                    self.result.grr_pct = float(pct_m)
                if pt_m is not None:
                    self.result.pt_ratio = float(pt_m)
                if ndc_m is not None:
                    self.result.ndc = float(ndc_m)
                # NDC stays at Python fallback (MSAVE doesn't include it).
            else:
                # Fallback: OCR the chart PNGs.
                logger.warning("GRR %s: worksheet print not parseable, "
                               "falling back to OCR", col)
                try:
                    augment_grr_with_ocr(
                        self.result, self.result.chart_paths or {},
                        fallback_pct=grr_pct_py,
                        fallback_pt=pt_ratio_py,
                        fallback_ndc=ndc_py)
                except Exception as ce:
                    logger.warning("OCR augment failed: %s; "
                                   "using Python fallback", ce)

        self.result.summary = (
            f"GRR% = {round(self.result.grr_pct,2) if self.result.grr_pct else 'N/A'}%  |  "
            f"P/T = {round(self.result.pt_ratio,4) if self.result.pt_ratio else 'N/A'}  |  "
            f"NDC = {round(self.result.ndc,1) if self.result.ndc else 'N/A'}  |  "
            f"Grade: {self.result.grade()}"
        )

        logger.info(f"GRR {col}: {self.result.summary}")
        return self.result

    def _run_minitab(self, mtb_path: str, df: pd.DataFrame, col: str, timeout:25, LED_SPECS) -> dict:
        """
        Generate GRR chart via Minitab using SET (no CSV files).
        Uses GageRR command with numeric Part/Operator columns.

        Workflow (mirrors ARR's approach):
          1. Build Minitab script with SET data + GageRR + XPPOINT
          2. Run Minitab, wait for exit (PowerPoint opens with the charts)
          3. Use win32com to GetActiveObject('Powerpoint.Application') and SaveAs PPTX
          4. Wait for images to fully load into the presentation
          5. Quit PowerPoint, then use python-pptx to extract images
          6. Cleanup: kill Minitab to clear workspace for next item
        """
        import subprocess as _sub
        import time as _time
        import gc as _gc
        out_dir = Path.home() / "grr_charts"
        out_dir.mkdir(exist_ok=True)
        ts = int(_time.time() * 1000)
        safe = col.replace("/", "_")
        mtb_file = out_dir / "grr_{}_{}.mtb".format(safe, ts)
        pptx_path = str(out_dir / "grr_{}_{}.pptx".format(safe, ts))
        img_dir = str(out_dir)
        grr_values_path = str(out_dir / "grr_{}_{}_vals.txt".format(safe, ts))

        # Kill existing Minitab / PowerPoint before starting
        for proc_name in ("Mtb.exe", "POWERPNT.EXE"):
            for _ in range(2):

                try:
                    _sub.run(["taskkill", "/F", "/IM", proc_name],
                             capture_output=True, timeout=5)
                except Exception:
                    pass

        for _ in range(10):
            result = _sub.getoutput('tasklist | find "POWERPNT.EXE"')
            if "POWERPNT.EXE" not in result:
                break
            # _sub.call("taskkill /f /im POWERPNT.EXE", shell=True)
            _sub.call("taskkill /f /im POWERPNT.EXE", shell=True)
            time.sleep(0.5)


        # ── Build GRR data from dataframe ─────────────────────────────────
        sn_col, op_col = self._detect_cols(df, col)
        df_work = df.copy()
        if sn_col:
            df_work["part_str"] = df_work[sn_col].astype(str).str.strip()
        else:
            df_work["part_str"] = df_work.index.astype(str)
        if op_col:
            df_work["op_str"] = df_work[op_col].astype(str).str.strip()
        else:
            df_work["op_str"] = "1"
        df_work = df_work.dropna(subset=["part_str", "op_str"])
        df_work = df_work[df_work["op_str"].str.lower().isin(["nan", "none", "parameter", ""]) == False]
        df_work = df_work[df_work["part_str"].str.lower().isin(["nan", "none", ""]) == False]
        unique_ops  = df_work["op_str"].unique().tolist()
        unique_parts = df_work["part_str"].unique().tolist()
        if len(unique_ops) == 0:
            logger.warning("_run_minitab: no operators in data, skipping")
            return {}
        logger.info("_run_minitab: unique_ops=%s n_parts=%d n_rows=%d",
                     unique_ops, len(unique_parts), len(df_work))
        logger.info("simon df_work is {}, col is {}".format(df_work, col))
        # led_col = col.lower() if col.lower() in df_work.columns else col.replace("_INTENSITY", "")
        # led_col = col.lower()
        led_col = col
        for i in df.columns:
            logger.info("simon i is {}".format(i))
        if led_col not in df_work.columns:
            logger.warning("_run_minitab: LED col '%s' not found, skipping", led_col)
            return {}
        df_work = df_work.dropna(subset=[led_col])
        logger.info("simon df_work is {}, led_col is {}".format(df_work, led_col))
        n_parts = len(unique_parts)
        n_ops   = len(unique_ops)

        # ── Build c1/c2/c3 in jason log time order (df's natural order) ──────
        # Keep df_work's original order (which is the upload/log time order).
        # C1 = part index cycling 1..N (e.g. 1, 2, … 10, 1, 2, … 10, …).
        #      This represents 10 SNs × 3 ops × 3 reps = 90 measurements
        #      (or however the actual study is structured). N defaults to 10
        #      from GRR_FORM_DEFAULTS['n_parts']; extra unique SNs in the data
        #      are ignored for C1 encoding.
        # C2 = inspector numbers split 30/30/30 (override QR_SCAN).
        # C3 = raw measurement value in the same order.
        n_expected_parts = int(GRR_FORM_DEFAULTS.get("n_parts", 10) or 10)
        n_parts_int = n_expected_parts
        logger.info("_run_minitab: %d unique parts in data, C1 cycles 1..%d",
                    len(unique_parts), n_parts_int)

        df_sorted = df_work.reset_index(drop=True)
        logger.info("simon df sorted is {}, led_col is ".format(df_sorted, led_col))
        # c3_vals = df_sorted[led_col].astype(int).tolist()
        logger.info(" simon df_sorted[led_col] is \n {} \n {}".format(df_sorted[led_col], df_sorted[led_col].astype(float)))
        c3_vals = df_sorted[led_col].astype(float).tolist()
        logger.info("simon c3_vals {}".format(c3_vals))
        n_rows  = len(c3_vals)
        # C1: clean repeating cycle 1, 2, … N, 1, 2, … N, …
        c1_vals = [((i % n_parts_int) + 1) for i in range(n_rows)]

        # SET C2: prefer inspector_numbers from GRR_FORM_DEFAULTS (split into
        # 3 groups of 30, one per inspector). Fall back to operator IDs from
        # the data when the config list is empty.
        c2_from_config = list(GRR_FORM_DEFAULTS.get("inspector_numbers") or [])
        if c2_from_config:
            # Pad to at least 3 entries (use last as fallback)
            while len(c2_from_config) < 3:
                c2_from_config.append(c2_from_config[-1] if c2_from_config else "1")
            n_per_block = max(1, n_rows // 3) if n_rows >= 3 else 1
            c2_vals = []
            for insp in c2_from_config[:3]:
                c2_vals.extend([str(insp)] * n_per_block)
            # Trim / pad to exactly n_rows
            if len(c2_vals) > n_rows:
                c2_vals = c2_vals[:n_rows]
            elif len(c2_vals) < n_rows:
                c2_vals.extend([c2_from_config[-1]] * (n_rows - len(c2_vals)))
        else:
            c2_vals = df_sorted["op_str"].astype(str).tolist()
        logger.info("_run_minitab: n_parts=%d n_ops=%d n_rows=%d (c2 source=%s)",
                    n_parts, n_ops, n_rows,
                    "config" if c2_from_config else "df")

        # LSL / USL from LED_SPECS for this item (V225 default).
        logger.info("simon run LED_SPECS is {}, col is {}".format(LED_SPECS, col))
        try:
            _spec = LED_SPECS.get(col) or LED_SPECS.get(_norm_led(col)) or {}
            lsl = float(_spec["lsl"]) if _spec.get("lsl") is not None else float(min(c3_vals))
            usl = float(_spec["usl"]) if _spec.get("usl") is not None else float(max(c3_vals))
        except Exception:
            lsl, usl = float(min(c3_vals)), float(max(c3_vals))
        logger.info("simon lsl and usl is {} {} \n {}".format(lsl, usl, _spec))
        # Subcommand values for Minitab's Gage/Date/User/GTolerance
        # metadata (Minitab 21+). Pulled from the analyzer input;
        # fall back to the column name / today / current user /
        # USL-LSL if the caller did not pass them in. Must be set
        # AFTER lsl/usl are computed above.
        import datetime as _dt
        _mtb_date  = getattr(self, "user_date", None) or _dt.date.today().strftime("%b %d %Y")
        # User comes from config `reported_by` (not from the OS env
        # variable, which would yield the build user on the
        # deployment machine).
        _mtb_user  = (getattr(self, "user_name", None)
                      or GRR_FORM_DEFAULTS.get("reported_by", "")
                      or "OpenClaw")
        _mtb_gtol  = getattr(self, "user_gtol", None) or round((usl - lsl), 3)

        # (LSL / USL already pulled from LED_SPECS above.)

        def fmt_row(vals, width=10, quote=False):
            rows = []
            for i in range(0, len(vals), width):
                row_vals = vals[i:i+width]
                if quote:
                    row_vals = [repr(str(v)) if " " in str(v) else str(v) for v in row_vals]
                rows.append(" ".join(str(v) for v in row_vals))
            return rows

        c1_rows = fmt_row(c1_vals)
        c2_rows = fmt_row(c2_vals, quote=True)
        c3_rows = fmt_row(c3_vals)

        lines = (
            [
                "NAME C1 'Part'",
                "NAME C2 'Operator'",
                "NAME C3 'Measurement'",
            ]
            + ["SET C1"]
            + c1_rows + ["END.", ""]
            + ["SET C2"]
            + c2_rows + ["END.", ""]
            + ["SET C3"]
            + c3_rows + ["END.", ""]
            + [
                "",
                # Main GageRR panel (R + Xbar + Components + plots).
                # NOTE: Minitab 19's GageRR command does NOT accept
                # the Gage/Date/User/GTolerance subcommands — those
                # are Minitab 21+ worksheet metadata fields, not
                # GageRR subcommands. Including them raises:
                #   "Subcommand GAGE is out of sequence or is not
                #    valid with this command."
                # Use the worksheet-level METADATA command BEFORE
                # GageRR if you need to tag the analysis with the
                # inspector / date / tolerance. For Minitab 19 we
                # omit them — the values still appear in the
                # OUTFILE session text and the report.
                #
                # OUTFILE captures the session window to a .txt
                # file so we can read the official Gage Evaluation
                # table values for %GR&R / %P/T / NDC. OUTFILE is a
                # session command, NOT a GageRR subcommand, and
                # Minitab 19 accepts it.
                "GageRR;",
                "  Parts 'Part';",
                "  Opers 'Operator';",
                "  Response 'Measurement';",
                "  Studyvar 6;",
                "  Pvalue 0.9999999;",
                "  LSL {};".format(lsl),
                "  USL {};".format(usl),
                "  Gage \"{}\";".format(col),
                "  Date \"{}\";".format(_mtb_date),
                "  User \"{}\";".format(_mtb_user),
                "  GTolerance \"{}\";".format(_mtb_gtol),
                "  Risk.",
                "",
                # Use JOURNAL to mirror the session output (which
                # includes the Gage Evaluation table) to a .txt
                # file. (DISABLED: Minitab 22 rejects JOURNAL with
                # "Unknown Minitab command: JOUR". OCR on the
                # PPTX-exported chart images is the working
                # alternative — see _ocr_extract_grr_metrics.)
                # 'JOURNAL "{}";'.format(grr_values_path.replace("\\", "\\\\")),
                "",
                "",
                "",
                "",
                "",
                "",
                "",
                "",
                "",
                "",
                "",
                "",
                # Send all current graphs to PowerPoint
                "XPPOINT.",
                "",
                # Clear Minitab workspace so the next item starts fresh
                # "NEW",
                "",
                "END",
            ]
        )

        for attempt in range(4):
            try:
                mtb_file.write_text(chr(10).join(lines), encoding="latin1")
                break
            except (PermissionError, OSError):
                if attempt < 3: _time.sleep(0.5)
                else: raise

        logger.info("GRR MTB written: %s", mtb_file)
        logger.info("PPTX target: %s", pptx_path)

        paths = {}
        worksheet_text_via_com = None
        try:
            # Run Minitab: GageRR creates charts, XPPOINT opens PowerPoint
            proc = _sub.Popen([mtb_path, str(mtb_file)],
                              stdout=_sub.PIPE, stderr=_sub.PIPE)
            try:
                logger.info('Minitab started (PID=%s)', proc.pid)
            except Exception:
                logger.info('Minitab started (PID unknown)')
            try:
                # _rc = proc.wait(timeout=30)
                _rc = proc.wait(timeout=timeout)
                logger.info('Minitab finished (exit code %s)', _rc)
            except _sub.TimeoutExpired:
                logger.warning('Minitab did not exit in 30s, terminating...')
                proc.terminate()
                try:
                    proc.wait(timeout=10)
                except Exception:
                    pass
            except Exception as we:
                logger.warning('Minitab wait failed: %s', we)

            # Pull the official %GR&R / %P/T / NDC from Minitab's
            # ActiveWorksheet columns M1–M30 via COM. This is the
            # only headless way that works on Minitab 22 — the
            # OUTFILE / MSAVE / JOURNAL session commands are all
            # rejected on Minitab 22 (raised "Subcommand GAGE is
            # out of sequence" or "Unknown Minitab command").
            # try:
            #     import win32com.client as _wc
            #     _mtb = _wc.GetActiveObject("Mtb.Application")
            #     _ws = _mtb.ActiveProject.ActiveWorksheet
            #     # Minitab puts GageRR Variability Components in
            #     # M1–M30 (5 source rows × 6 columns: VarComp,
            #     # %VarComp, StDev, %SV, 6*SD, SV/Toler).
            #     _nums = []
            #     for _col_idx in range(1, 31):
            #         try:
            #             _col = _ws.Columns(_col_idx)
            #             # Get the M1..M30 column data. Each column
            #             # has 5 rows for the 5 source categories.
            #             for _row in range(1, 6):
            #                 try:
            #                     _v = _col.Cells(_row, 1).Value
            #                     if _v is not None:
            #                         _nums.append(float(_v))
            #                 except Exception:
            #                     pass
            #         except Exception:
            #             pass
            #     # Build a mock Minitab session-style text so the
            #     # existing _parse_minitab_session_for_grr_str
            #     # parser can read it.
            #     if _nums:
            #         _lines = ["Source  VarComp  %VarComp  StDev  %SV  6*SD  SV/Toler"]
            #         _source_names = [
            #             "Total Gage R&R", "Repeatability",
            #             "Reproducibility", "Part-to-Part", "Total Variation"]
            #         for _i, _name in enumerate(_source_names):
            #             _chunk = _nums[_i*6:(_i+1)*6]
            #             if len(_chunk) >= 4:
            #                 _line = "{}  {}  {}  {}  {}  {}  {}".format(
            #                     _name, *_chunk[:6])
            #                 _lines.append(_line)
            #         worksheet_text_via_com = "\n".join(_lines)
            #         logger.info("COM: captured %d M1-M30 values from "
            #                     "Minitab worksheet", len(_nums))
            # except Exception as ce:
            #     logger.warning("COM read from Minitab failed: %s", ce)
            #
            # # Give PowerPoint a moment to fully load
            # logger.info('Waiting 5s for PowerPoint to initialise...')
            # _time.sleep(5)

            # Save PowerPoint via COM
            try:
                import win32com.client
                import os as _os
                _os.makedirs(out_dir, exist_ok=True)

                logger.info('Connecting to open PowerPoint via GetActiveObject...')
                ppt_app = win32com.client.GetActiveObject('Powerpoint.Application')
                presentation = ppt_app.ActivePresentation
                # ppSaveAsOpenXMLPresentation = 24
                presentation.SaveAs(pptx_path, FileFormat=24)
                logger.info('PPTX saved via COM: %s', pptx_path)

                # Wait for images to load
                logger.info('Waiting 15s for all images to load into PPTX...')
                # _time.sleep(3)

                try:
                    presentation.Saved = True
                    ppt_app.Quit()
                    logger.info('PowerPoint closed')
                except Exception as e:
                    logger.warning('PowerPoint.Quit failed: %s', e)

                del presentation
                del ppt_app
                _gc.collect()
                _time.sleep(2.0)
            except ImportError:
                logger.warning('win32com not available, PPTX save skipped')
            except Exception as e:
                logger.warning('PowerPoint COM SaveAs failed: %s', e)

            # Wait for PPTX file to be released
            for _w in range(6):
                if Path(pptx_path).exists():
                    break
                logger.info('Waiting for PPTX to be released (%d)...', _w)
                _time.sleep(1.0)

            # Extract images from PPTX
            if Path(pptx_path).exists():
                logger.info('PPTX confirmed: %s (size=%d)',
                            pptx_path, Path(pptx_path).stat().st_size)
                try:
                    from pptx import Presentation
                    from pptx.enum.shapes import MSO_SHAPE_TYPE

                    prs = Presentation(pptx_path)
                    ext_map = {'image/jpeg': '.jpg', 'image/png': '.png',
                               'image/gif': '.gif', 'image/bmp': '.bmp'}

                    extracted = []
                    for slide_idx, slide in enumerate(prs.slides, 1):
                        for shape_idx, shape in enumerate(slide.shapes, 1):
                            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                                continue
                            img = shape.image
                            ext = ext_map.get(img.ext, '.png')
                            out_name = 'slide{:02d}_shape{:02d}{}'.format(
                                slide_idx, shape_idx, ext)
                            out_path = _os.path.join(img_dir, out_name)
                            with open(out_path, 'wb') as f:
                                f.write(img.blob)
                            extracted.append(out_path)
                            logger.info('Extracted: %s', out_path)

                    # Map by stem so Excel generator can pick chart by name
                    for p in extracted:
                        paths[Path(p).stem] = p

                    logger.info('GRR charts extracted: %d', len(extracted))

# Cleanup PPTX
                    try:
                        Path(pptx_path).unlink()
                        logger.info('PPTX deleted')
                    except Exception as e:
                        logger.warning('Could not delete PPTX: %s', e)
                except ImportError:
                    logger.warning('python-pptx not installed')
                except Exception as e:
                    logger.warning('PPTX extraction failed: %s', e)
            else:
                logger.warning('PPTX not found: %s', pptx_path)

        except Exception as e:
            logger.warning("Minitab GRR failed: %s", e)

        finally:
            # Always kill Minitab to clear workspace for next item
            try:
                _sub.run(["taskkill", "/F", "/IM", "Mtb.exe"],
                         capture_output=True, timeout=5)
            except Exception:
                pass

        paths['_img_dir'] = img_dir
        if worksheet_text_via_com:
            paths['_grr_worksheet'] = worksheet_text_via_com
        # Belt-and-suspenders: try the disk worksheet file too
        # (V245 JOURNAL path).
        try:
            if Path(grr_values_path).exists():
                if "_grr_worksheet" not in paths:
                    paths['_grr_worksheet'] = Path(grr_values_path).read_text(
                        encoding="latin1", errors="ignore")
                Path(grr_values_path).unlink(missing_ok=True)
        except Exception as se:
            logger.warning("Could not capture worksheet print %s: %s",
                           grr_values_path, se)
        return paths

    @staticmethod
    def _detect_cols_static(df: pd.DataFrame, col: str):
        """
        Detect part-number and operator column names from the GRR DataFrame.

        Priority (first match wins, scanning in priority order):
          Part column:
            'part_num'  → numeric 1-11 from assign_trials (intermediate Excel)
            'sample'    → GRR template (1-10)
            'sn'        → device serial (MT7 raw log, fallback only)
          Operator column:
            'appraiser' → intermediate Excel / MT7 QR_SCAN result
            'inspector' → GRR template
            'qr_scan'   → MT7 raw QR_SCAN field
        """
        # Part column: part_num > sample > sn
        sn_col = None
        for candidate in ["part_num", "sample", "sn"]:
            for c in df.columns:
                if c.lower() == candidate:
                    sn_col = c
                    break
            if sn_col:
                break

        # Operator column: appraiser > inspector > qr_scan
        op_col = None
        for candidate in ["appraiser", "inspector", "qr_scan"]:
            for c in df.columns:
                if c.lower() == candidate:
                    op_col = c
                    break
            if op_col:
                break

        return sn_col, op_col

    @staticmethod
    def run_all_minitab(df: pd.DataFrame, items: list, minitab_path: str, LED_SPECS) -> dict:
        """
        Run GRR GageRR for all items, ONE Minitab session per item.

        Each iteration:
          1. Build a per-item Minitab script with SET data + GageRR + XPPOINT
          2. Run Minitab, wait for exit (PowerPoint opens)
          3. Use win32com to GetActiveObject('Powerpoint.Application') and SaveAs PPTX
          4. Wait for images to load, then Quit PowerPoint
          5. Use python-pptx to extract images from PPTX
          6. Kill Minitab to clear workspace (prevents next item from inheriting charts)

        Returns a flat dict mapping "<safe>__<slideN>_shape<N>" → PNG path.
        Per-item chart groups are kept separate by prefixing keys with the
        item's safe name to avoid cross-item collisions.

        See GRRAnalyzer._run_minitab (single-item) for the detailed workflow.
        """
        import subprocess as _sub
        import time as _time

        all_paths = {}
        last_img_dir = ""
        for idx, item in enumerate(items):
            if idx==0:
                timeout=40
            else:
                timeout=33
            norm = _norm_led(item)
            if norm not in df.columns and item not in df.columns:
                logger.warning("run_all_minitab: column '%s' not in df, skipping", item)
                continue

            logger.info("=" * 60)
            logger.info("GRR per-item: %d/%d  %s", idx + 1, len(items), item)
            logger.info("=" * 60)

            # Build a single-item analyzer and call _run_minitab directly
            analyzer = GRRAnalyzer(df, item, LED_SPECS)
            item_paths = analyzer._run_minitab(minitab_path, df, norm, timeout, LED_SPECS)

            if not item_paths:
                logger.warning("run_all_minitab: no paths returned for '%s'", item)
                continue

            # Capture the per-item img_dir so we can wipe it after the
            # bytes have been moved to memory.
            this_img_dir = item_paths.get("_img_dir", "")
            ppt_candidate = this_img_dir and Path(this_img_dir).parent
            ppt_files_before = (
                list(ppt_candidate.glob("grr_*.pptx"))
                if ppt_candidate and ppt_candidate.exists() else []
            )
            # Prefix keys with the item's safe name to disambiguate across
            # items. Read the PNG bytes into BytesIO then delete the file
            # immediately so we don't accumulate disk artefacts across
            # the batch — the Excel generator will consume the in-memory
            # BytesIO objects when it embeds the charts.
            from io import BytesIO as _BytesIO
            safe = item.replace("/", "_")
            pngs_deleted = 0
            for k, v in item_paths.items():
                if k == "_img_dir":
                    continue
                p = Path(v)
                try:
                    if p.exists() and p.is_file():
                        with open(p, "rb") as fh:
                            buf = _BytesIO(fh.read())
                        buf.name = p.name   # openpyxl uses this for extension
                        all_paths["{}__".format(safe) + k] = buf
                        try:
                            p.unlink()
                            pngs_deleted += 1
                        except Exception as de:
                            logger.warning("Could not delete PNG %s: %s", p, de)
                    else:
                        logger.warning("PNG missing for %s/%s: %s", item, k, p)
                except Exception as le:
                    logger.warning("Could not load %s to memory: %s", p, le)
            logger.info("Cleaned %d PNG files for item %s", pngs_deleted, item)

            # Belt-and-suspenders: explicitly delete any pptx files for
            # this item (in case _run_minitab's own cleanup failed).
            ppt_deleted = 0
            for pp in ppt_files_before:
                try:
                    if pp.exists():
                        pp.unlink()
                        ppt_deleted += 1
                        logger.info("Deleted PPTX: %s", pp)
                except Exception as pe:
                    logger.warning("Could not delete PPTX %s: %s", pp, pe)
            logger.info("Cleaned %d PPTX files for item %s", ppt_deleted, item)

            # Try to clean the empty per-item img_dir.
            if this_img_dir:
                try:
                    _p = Path(this_img_dir)
                    if _p.exists() and _p.is_dir():
                        for leftover in _p.iterdir():
                            try:
                                leftover.unlink()
                                logger.info("Deleted leftover: %s", leftover)
                            except Exception:
                                pass
                        if not any(_p.iterdir()):
                            _p.rmdir()
                            logger.info("Removed empty img_dir: %s", _p)
                except Exception as rde:
                    logger.warning("Could not rmdir %s: %s", this_img_dir, rde)
                last_img_dir = this_img_dir

            # Brief pause between items
            _time.sleep(1.0)

        # Pass the img_dir through to the Excel generator (still useful
        # for diagnostics; in-memory chart paths are in all_paths).
        if last_img_dir:
            all_paths["_img_dir"] = last_img_dir
            # Save the MSAVE / PRINT worksheet text for each item so
            # compute() can later pull the official Minitab %GR&R /
            # %P/T values without OCR.
            try:
                _vp = ppt_candidate / Path(grr_values_path).name
                if _vp and _vp.exists():
                    all_paths["{}__grr_worksheet".format(safe)] = _vp.read_text(
                        encoding="latin1", errors="ignore")
                    _vp.unlink(missing_ok=True)
            except Exception as se:
                logger.warning("Could not capture worksheet print: %s", se)

        logger.info("=" * 60)
        logger.info("GRR per-item total: %d charts across %d items (all in-memory)",
                    len([k for k in all_paths if not k.startswith("_")]),
                    len(items))
        logger.info("=" * 60)

        return all_paths
